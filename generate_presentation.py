import os
import io
import pandas as pd
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def get_thumbnail(video_id):
    """
    Intenta descargar el thumbnail en alta resolución (maxresdefault).
    Si falla (404), intenta con la resolución estándar (hqdefault).
    """
    urls_to_try = [
        f"https://img.youtube.com/vi/{video_id}/maxresdefault.jpg",
        f"https://img.youtube.com/vi/{video_id}/hqdefault.jpg"
    ]
    for url in urls_to_try:
        response = requests.get(url)
        if response.status_code == 200:
            return io.BytesIO(response.content)
    return None

def generate_pptx():
    csv_file = 'top_videos_report.csv'
    output_file = 'Presentacion_Top_Videos.pptx'

    if not os.path.exists(csv_file):
        print(f"El archivo {csv_file} no existe. Por favor, corre main.py o top_video_report.py primero.")
        return

    # Leer el CSV
    df = pd.read_csv(csv_file)
    if df.empty:
        print("El archivo CSV está vacío.")
        return

    # Crear la presentación
    prs = Presentation()
    
    # --- Diapositiva de Título ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 1. Fondo de la portada
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x75, 0x9c, 0xcf)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "TOP VIDEOS DE YOUTUBE"
    
    # Formatear el título de la presentación
    p_title = title.text_frame.paragraphs[0]
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    
    # Remover el subtítulo para que solo quede el título principal
    subtitle.text = ""


    # Agrupar por canal
    grouped = df.groupby('ChannelName')

    for channel_name, group in grouped:
        for i in range(0, len(group), 2):
            chunk = group.iloc[i:i+2]
            
            # Crear diapositiva en blanco para los videos (layout 6)
            blank_slide_layout = prs.slide_layouts[6]
            video_slide = prs.slides.add_slide(blank_slide_layout)
            
            # 1. Agregar color de fondo #759ccf
            background = video_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x75, 0x9c, 0xcf)
            
            # 2. Título del canal arriba a la izquierda
            txBox_channel = video_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            tf_channel = txBox_channel.text_frame
            p_channel = tf_channel.add_paragraph()
            p_channel.text = str(channel_name).upper()
            p_channel.font.bold = True
            p_channel.font.size = Pt(28)
            p_channel.font.color.rgb = RGBColor(255, 255, 255)
            
            # Alturas base para Top y Bottom video
            base_tops = [Inches(1.2), Inches(4.3)]
            
            for j, (index, row) in enumerate(chunk.iterrows()):
                base_top = base_tops[j]
                
                # A. Descargar e incrustar la imagen del video (Izquierda)
                video_id = row['VideoID']
                img_stream = get_thumbnail(video_id)
                if img_stream:
                    img_width = Inches(4)
                    left_img = Inches(0.5)
                    video_slide.shapes.add_picture(img_stream, left_img, base_top, width=img_width)
                
                # B. Título del video (Derecha)
                txBox_title = video_slide.shapes.add_textbox(Inches(4.8), base_top, Inches(4.7), Inches(1))
                tf_title = txBox_title.text_frame
                tf_title.word_wrap = True
                p_title = tf_title.paragraphs[0]
                p_title.text = f"+ {str(row['Title']).upper()}"
                p_title.font.bold = True
                p_title.font.size = Pt(16)
                p_title.font.color.rgb = RGBColor(255, 255, 255)
                
                # C. Datos del video (Debajo del título a la derecha)
                txBox_data = video_slide.shapes.add_textbox(Inches(4.8), base_top + Inches(1.2), Inches(4.7), Inches(1))
                tf_data = txBox_data.text_frame
                
                p_format = tf_data.paragraphs[0]
                p_format.text = f"FORMATO: {row['VideoType']}"
                p_format.font.size = Pt(14)
                p_format.font.bold = True
                p_format.font.color.rgb = RGBColor(255, 255, 255)
                
                p_views = tf_data.add_paragraph()
                p_views.text = f"VISUALIZACIONES: {row['ViewCount']:,}"
                p_views.font.size = Pt(14)
                p_views.font.bold = True
                p_views.font.color.rgb = RGBColor(255, 255, 255)
            
    # Guardar presentación
    prs.save(output_file)
    print(f"¡Éxito! Presentación guardada en '{output_file}'")

if __name__ == "__main__":
    generate_pptx()
